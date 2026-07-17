"""Generate the SIMPLIFIED Week 2 progress presentation.

Per supervisor feedback: minimal theory/formulas, most of the slot spent
live-demoing the dashboard on one machine. This deck is the short
opening + closing bookends around that live demo, not a methodology deep-dive.

Usage (from project root):
    python -m scripts.generate_week2_simple_pptx

Produces docs/Week2_Progress_Presentation_Simple.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs" / "Week2_Progress_Presentation_Simple.pptx"
SCRATCH = Path(r"C:\Users\lenovo\AppData\Local\Temp\claude\C--Documents-COMPETITION-Astra-Case2\790afb06-41ff-4409-a6d7-72d15e2afcbf\scratchpad")

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x1F, 0x6F, 0xC4)
GOOD = RGBColor(0x1E, 0x7B, 0x34)
BAD = RGBColor(0xB0, 0x2A, 0x2A)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)
    return shp


def add_topbar(slide, kicker, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background(); bar.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(11.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = kicker.upper()
    r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0); r.font.name = "Calibri"

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.42), Inches(12), Inches(0.65))
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(30); r2.font.bold = True
    r2.font.color.rgb = WHITE; r2.font.name = "Calibri"

    foot = slide.shapes.add_textbox(Inches(11.6), Inches(7.12), Inches(1.5), Inches(0.3))
    fp = foot.text_frame.paragraphs[0]; fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run(); fr.text = "Case 2 · WINTEQ"
    fr.font.size = Pt(9); fr.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.7, top=1.55, width=11.9, height=5.3,
                 font_size=20, color=RGBColor(0x22, 0x27, 0x2E), line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        text, level = item[0], item[1]
        col = item[2] if len(item) > 2 and item[2] else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level; p.space_after = Pt(12); p.line_spacing = line_spacing
        bullet = {0: "■  ", 1: "–  "}.get(level, "•  ")
        r = p.add_run()
        r.text = (bullet if level < 2 else "   " + bullet) + text
        r.font.size = Pt(font_size - level * 3)
        r.font.color.rgb = col; r.font.name = "Calibri"
        if level == 0:
            r.font.bold = True
    return tb


def add_note(slide, text, top=6.85, color=GRAY, size=12, italic=True):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.italic = italic; r.font.color.rgb = color


def add_image(slide, path, left, top, width=None, height=None, caption=None):
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    width=Inches(width) if width else None,
                                    height=Inches(height) if height else None)
    pic.line.color.rgb = RGBColor(0xD5, 0xDA, 0xE1); pic.line.width = Pt(0.75)
    if caption:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top + (height or 3) + 0.05),
                                       Inches(width or 5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GRAY
    return pic


def add_bignum(slide, left, top, width, number, label, color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(0xD5, 0xDA, 0xE1)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = number
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(14); r2.font.color.rgb = GRAY


# ---------------------------------------------------------------------
# Slide 1 - Title
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "CASE 2 · AOP WINTEQ BOOTCAMP · WEEK 2"
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)

tb = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Predictive Maintenance Dashboard"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tb.text_frame.add_paragraph()
r2 = p2.add_run(); r2.text = "Live Prototype Demo"
r2.font.size = Pt(44); r2.font.bold = True; r2.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(10.8), Inches(0.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "One machine. Watch the alarm and the failure prediction happen live."
r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = RGBColor(0xD5, 0xE3, 0xF2)

tb = s.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(10), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Track 5 — Dashboard · July 2026"
r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)

# ---------------------------------------------------------------------
# Slide 2 - Problem & goal
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "The Problem", "Why Build This?")
add_bullets(s, [
    ("Today: maintenance finds out something is wrong only after it is already visible on the machine.", 0),
    ("Single-sensor alarms are noisy — operators start ignoring them, including the real ones.", 0),
    ("What we want instead:", 0),
    ("Catch problems early, before they cause downtime.", 1),
    ("Fewer false alarms — trust the alarm again.", 1),
    ("An estimate of how much time is left before failure.", 1),
    ("One simple screen that shows all of this at a glance.", 1),
])

# ---------------------------------------------------------------------
# Slide 3 - Our data
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Our Data", "One Motor Bearing, Monitored to Failure")
add_bullets(s, [
    ("Astra's real sensor data is not ready yet, so we are proving the concept first on a public dataset built for exactly this problem: a bearing run from healthy to failure.", 0),
    ("Real vibration sensor readings, recorded every 10 minutes, for the machine's entire life.", 0),
    ("Honest note: today's prototype uses vibration only. Real deployment will add temperature, current and pressure — same dashboard, more sensors, once Astra's data arrives.", 0),
])

# ---------------------------------------------------------------------
# Slide 4 - How it works (plain language, no formulas)
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "How It Works", "In Plain Language")

boxes = [
    ("SENSOR\nREADING", ACCENT, 0.7),
    ("IS THIS\nNORMAL?", ACCENT, 3.5),
    ("HOW MUCH TIME\nIS LEFT?", GOOD, 6.3),
    ("DASHBOARD", GOOD, 9.1),
]
for label, color, left in boxes:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(2.3), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.color.rgb = NAVY
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
arrow_y = 2.55
for x in (3.05, 5.85, 8.65):
    arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(arrow_y), Inches(0.4), Inches(0.4))
    arr.fill.solid(); arr.fill.fore_color.rgb = GRAY; arr.line.fill.background()

add_bullets(s, [
    ("We compare each new reading to what “healthy” looked like for this machine.", 0),
    ("Once several readings look abnormal together — not just one noisy blip — the system raises an alarm.", 0),
    ("Only after the alarm fires does it start estimating a countdown to failure — it never guesses before there is real evidence.", 0),
    ("Everything above is shown live on one screen: a health score, an alarm status, and a time-remaining estimate.", 0),
], top=3.7, height=3, font_size=18)

# ---------------------------------------------------------------------
# Slide 5 - What to watch for (live demo intro)
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.9))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Live Demo — One Machine, Start to Failure"
r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE

add_bullets(s, [
    ("We will replay one real motor bearing, minute by minute, from perfectly healthy to failure.", 0),
    ("Watch three things on screen:", 0),
    ("Health score — starts at 100, watch it decline.", 1),
    ("Alarm status — Normal → Warning → Critical, and why it changed.", 1),
    ("Remaining time estimate — appears once the alarm fires, then counts down.", 1),
], left=0.9, top=2.0, width=11.5, height=4,
   color=RGBColor(0xE9, 0xF1, 0xFB))
add_note(s, "[Switch to the live dashboard now]", color=RGBColor(0x9C, 0xC5, 0xF0), italic=False, size=14)

# ---------------------------------------------------------------------
# Slide 6 - This week: found and fixed a bug
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "This Week", "We Found & Fixed a Real Bug")
add_bullets(s, [
    ("While preparing this demo, our numbers looked wrong on our longest test run — so we investigated instead of ignoring it.", 0),
    ("What we found: a single noisy sensor reading could trigger a false alarm — while the machine was still 100% healthy.", 0),
    ("That false alarm was throwing off every prediction that came after it.", 0),
    ("Fix: an alarm now needs two consecutive abnormal readings before it counts as real, not just one.", 0),
])
add_image(s, SCRATCH / "healthscore_trigger_comparison.png", left=1.6, top=4.55, width=10.1, height=2.55)

# ---------------------------------------------------------------------
# Slide 7 - Honest limitations
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Transparency", "What We're Still Working On")
add_bullets(s, [
    ("Vibration only for now — temperature, current and pressure are next, once real Astra data is available.", 0),
    ("The time-remaining estimate can still swing in the first hours after an alarm fires, and settles as more data comes in.", 0),
    ("Tested on one operating condition (constant speed/load). Real factory motors vary — handling that is on the roadmap.", 0),
])

# ---------------------------------------------------------------------
# Slide 8 - Next steps
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Roadmap", "Next Steps (Week 3)")
add_bullets(s, [
    ("Improve the stability of the time-remaining prediction right after an alarm fires.", 0),
    ("Prepare the dashboard to accept new sensor types (temperature, current, pressure) once Astra's real data is ready.", 0),
    ("Check the same approach on other machines in the dataset, not just this one.", 0),
    ("Start designing for live, real-time data instead of replaying historical files.", 0),
])

# ---------------------------------------------------------------------
# Slide 9 - Thank you
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11), Inches(1))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Thank You"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(11), Inches(0.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Questions & Discussion"
r.font.size = Pt(20); r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)

OUT_PATH.parent.mkdir(exist_ok=True)
prs.save(str(OUT_PATH))
print("Saved:", OUT_PATH)
print("Slide count:", len(prs.slides._sldIdLst))
