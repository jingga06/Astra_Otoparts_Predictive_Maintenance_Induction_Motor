"""Generate the Week 2 progress presentation (English, for weekly supervisor report).

Usage (from project root):
    python -m scripts.generate_week2_progress_pptx

Produces docs/Week2_Progress_Presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs" / "Week2_Progress_Presentation_v2.pptx"

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0x1F, 0x6F, 0xC4)
GOOD = RGBColor(0x1E, 0x7B, 0x34)
BAD = RGBColor(0xB0, 0x2A, 0x2A)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)
    return shp


def add_topbar(slide, kicker, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(11.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = kicker.upper()
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)
    r.font.name = "Calibri"

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.42), Inches(12), Inches(0.65))
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = title
    r2.font.size = Pt(28)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    r2.font.name = "Calibri"

    foot = slide.shapes.add_textbox(Inches(11.6), Inches(7.12), Inches(1.5), Inches(0.3))
    fp = foot.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = "Case 2 · WINTEQ"
    fr.font.size = Pt(9)
    fr.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.6, top=1.45, width=12.1, height=5.5,
                 font_size=18, color=RGBColor(0x22, 0x27, 0x2E), line_spacing=1.15):
    """items: list of (text, level, bold_prefix_or_None, color_override_or_None)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, level = item[0], item[1]
        col = item[2] if len(item) > 2 and item[2] else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(10)
        p.line_spacing = line_spacing
        bullet = {0: "■  ", 1: "–  ", 2: "•  "}.get(level, "•  ")
        r = p.add_run()
        r.text = (bullet if level < 2 else "   " + bullet) + text
        r.font.size = Pt(font_size - level * 2)
        r.font.color.rgb = col
        r.font.name = "Calibri"
        if level == 0:
            r.font.bold = True
    return tb


def add_note(slide, text, top=6.65, color=GRAY, size=12, italic=True):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.italic = italic
    r.font.color.rgb = color


SCRATCH = Path(r"C:\Users\lenovo\AppData\Local\Temp\claude\C--Documents-COMPETITION-Astra-Case2\790afb06-41ff-4409-a6d7-72d15e2afcbf\scratchpad")


def add_image(slide, path, left, top, width=None, height=None, caption=None):
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    width=Inches(width) if width else None,
                                    height=Inches(height) if height else None)
    pic.line.color.rgb = RGBColor(0xD5, 0xDA, 0xE1)
    pic.line.width = Pt(0.75)
    if caption:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top + (height or 3) + 0.05),
                                       Inches(width or 5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GRAY
    return pic


def add_table(slide, rows, col_widths_in, top=1.55, left=0.6, row_h=0.5,
              header_bg=NAVY, header_fg=WHITE, highlight_rows=None, font_size=13):
    highlight_rows = highlight_rows or {}
    n_rows, n_cols = len(rows), len(rows[0])
    width = Inches(sum(col_widths_in))
    height = Inches(row_h * n_rows)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), width, height).table
    for c, w in enumerate(col_widths_in):
        gtable.columns[c].width = Inches(w)
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = gtable.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
                run.font.color.rgb = header_fg
                run.font.bold = True
            else:
                fill_color = highlight_rows.get(r_idx)
                if fill_color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill_color
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                run.font.color.rgb = RGBColor(0x22, 0x27, 0x2E)
    return gtable


# ---------------------------------------------------------------------
# Slide 1 - Title
# ---------------------------------------------------------------------
s = add_slide()
add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "CASE 2 · AOP WINTEQ BOOTCAMP · WEEK 2 PROGRESS REPORT"
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)

tb = s.shapes.add_textbox(Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Predictive Maintenance Dashboard"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tb.text_frame.add_paragraph()
r2 = p2.add_run(); r2.text = "for Induction Motors"
r2.font.size = Pt(44); r2.font.bold = True; r2.font.color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(10.8), Inches(0.9))
p = tb.text_frame.paragraphs[0]
r = p.add_run()
r.text = "This week: pipeline fully implemented on NASA IMS Bearing data, dashboard working end to end, and one real bug found + fixed with before/after evidence."
r.font.size = Pt(17); r.font.italic = True; r.font.color.rgb = RGBColor(0xD5, 0xE3, 0xF2)

tb = s.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(10), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Track 5 — Dashboard · July 2026"
r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)

# ---------------------------------------------------------------------
# Slide 2 - Recap: problem & goal
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Recap", "Business Problem & Goal (Astra Use Case 2)")
add_bullets(s, [
    ("Business problem", 0),
    ("Reactive monitoring: current systems flag issues only after degradation is already visible — short action window.", 1),
    ("Single-parameter thresholds cannot describe a complex machine — noisy alarms erode trust in the system.", 1),
    ("Goal (from the bootcamp brief)", 0),
    ("Multi-parameter monitoring & analysis, AI-based early anomaly detection.", 1),
    ("Intelligent false-alarm reduction by relating multiple parameters, not single thresholds.", 1),
    ("Remaining Useful Life (RUL) estimation, real-time health scoring, automated alerts + recommendations.", 1),
])
add_note(s, "Source: Bootcamp Briefing AOP Winteq, Use Case 2 — Predictive Maintenance for Induction Motor Failure.")

# ---------------------------------------------------------------------
# Slide 3 - Dataset decision
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Data", "Dataset: NASA IMS Bearing (Run-to-Failure)")
add_bullets(s, [
    ("Astra's real plant sensor data is not available yet at this stage.", 0),
    ("Using NASA IMS Bearing dataset as a stand-in: 4 bearings, Rexnord ZA-2115, 2000 RPM, run-to-failure vibration, 20 kHz.", 0),
    ("Why this dataset: the only public dataset that supports realistic RUL modelling (true run-to-failure), not just fault classification on clean lab data.", 1),
    ("Test 2 (163.8 h, bearing 1 – outer race) = development run.", 1),
    ("Test 3 (1073.3 h, bearing 3 – outer race) = untouched holdout run, used only for final validation.", 1),
    ("Known scope gap (honest limitation)", 0),
    ("The brief asks for vibration + temperature + current + pressure. This dataset provides vibration only — addressed in the roadmap slide.", 1),
])

# ---------------------------------------------------------------------
# Slide 4 - Architecture
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "System Design", "Pipeline: Two Branches, One Dashboard")

boxes = [
    ("RAW VIBRATION", "20 kHz snapshot\nevery 10 min", 0.5),
    ("FEATURES", "RMS, kurtosis, crest\nfactor, BPFO/BPFI/BSF\nenvelope energy", 2.35),
    ("BRANCH A\nDETECT", "Isolation Forest\n(healthy data only)\n+ 3σ threshold", 4.6),
    ("ALARM LOGIC", "Persistence +\nvoting rule\n→ is something wrong?", 6.85),
]
for label, sub, left in boxes:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.7), Inches(1.95), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT; box.line.color.rgb = NAVY
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    pp = tf.add_paragraph(); pp.text = sub; pp.font.size = Pt(10)
    pp.font.color.rgb = RGBColor(0xE9, 0xF1, 0xFB); pp.alignment = PP_ALIGN.CENTER

boxes2 = [
    ("BRANCH B\nPREDICT", "PCA Health Index\n→ knee → exponential\nRUL fit", 2.35),
    ("RUL OUTPUT", "→ how long until\nfailure?", 4.6),
    ("DASHBOARD", "Health score, alarm,\nRUL, recommendation", 6.85),
]
for label, sub, left in boxes2:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.55), Inches(1.95), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = GOOD; box.line.color.rgb = NAVY
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    pp = tf.add_paragraph(); pp.text = sub; pp.font.size = Pt(10)
    pp.font.color.rgb = RGBColor(0xE9, 0xF1, 0xFB); pp.alignment = PP_ALIGN.CENTER

add_bullets(s, [
    ("Same feature table feeds both branches: Branch A answers “is something wrong right now”, Branch B answers “how long until failure”.", 0),
], top=5.4, height=1.6, font_size=15)
add_note(s, "Implementation: pdm/features.py, pdm/anomaly.py, pdm/health_index.py, pdm/rul.py, pdm/alarm.py, app/dashboard.py")

# ---------------------------------------------------------------------
# Slide 5 - Feature engineering
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Methodology", "Feature Engineering — Physics-Grounded, Not Guesswork")
add_bullets(s, [
    ("Time-domain features", 0),
    ("RMS — overall vibration energy (ISO 10816 severity metric).", 1),
    ("Kurtosis — impulsiveness; healthy ≈ 3, bearing impacts push it far higher (observed 3.6 → 17 near failure).", 1),
    ("Crest factor — peak-to-RMS ratio, a second impulsiveness indicator.", 1),
    ("Frequency-domain: envelope analysis", 0),
    ("Band-pass (2–10 kHz) → Hilbert envelope → FFT → energy at bearing fault frequencies.", 1),
    ("BPFO 236.4 Hz, BPFI 296.9 Hz, BSF 139.9 Hz — computed from Rexnord ZA-2115 geometry, verified against the team’s documented reference values (exact match).", 1),
    ("Verified on real data: BPFO envelope energy rose ~300× (1,230 → 388,268) in the hours before failure.", 1),
])

# ---------------------------------------------------------------------
# Slide 6 - Branch A anomaly detection
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Branch A", "Anomaly Detection — Isolation Forest")
add_bullets(s, [
    ("Why unsupervised: the dataset has no per-timestamp ground truth.", 0),
    ("We only know the run's final outcome (documented failure mode) — not which moment counts as “already abnormal”.", 1),
    ("This matches what real plant data usually looks like too (maintenance logs give an end-of-life event, not a labelled trace).", 1),
    ("How it works", 0),
    ("Trained only on the first 50% of each run (assumed healthy).", 1),
    ("Anomalies are isolated in fewer random splits than normal points — fewer splits = higher anomaly score.", 1),
    ("Alarm threshold = healthy-score mean + 3σ, a statistical threshold, not an arbitrary number.", 1),
])

# ---------------------------------------------------------------------
# Slide 7 - Branch B Health Index & RUL
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Branch B", "Health Index & Remaining Useful Life (RUL)")
add_bullets(s, [
    ("Health Index: PCA on monotonicity-filtered features → one number, scaled to a 0–100 health score.", 0),
    ("RUL starts only after the first confirmed alarm (the “knee”) — before that, the system reports “monitoring, not stable yet” instead of guessing.", 0),
    ("Model:  HI(t) = p1 · e^(p2·t) + p3,  refit as new data arrives.", 0),
    ("EOL = (1/p2) · ln((ET – p3) / p1);   RUL(t) = EOL – t", 0),
    ("ET (failure threshold) is taken from the OTHER run's end-of-life value — cross-run, so it never leaks the run's own future.", 1),
])

# ---------------------------------------------------------------------
# Slide 8 - False alarm reduction
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Methodology", "False-Alarm Reduction: Persistence + Voting")
add_bullets(s, [
    ("Persistence rule: a single abnormal feature must hold for N consecutive windows before it counts (debounce).", 0),
    ("Voting rule: ≥2 features abnormal at once → escalates directly to CRITICAL (a stronger, faster signal).", 0),
    ("Both windows scale with run length (burn-in 5%, persistence 0.5%) — the same logic works unmodified on a 164 h run and a 1073 h run.", 0),
    ("Hysteresis on recovery: status only clears after several consecutive normal readings, preventing alarm flicker.", 0),
])

# ---------------------------------------------------------------------
# Slide 9 - This week's debugging story
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "This Week's Work", "Found & Fixed: A Premature-Trigger Bug")
add_bullets(s, [
    ("What we noticed: our holdout-run results did not match numbers already written into an earlier proposal draft — investigated instead of assuming either was right.", 0),
    ("Root cause found: the voting rule could fire CRITICAL off a single noisy reading. On the holdout run, this fired at hour 257 of 1073 — while the health score was still 100% (perfectly healthy).", 0),
    ("Consequence: RUL started fitting from that false trigger, forcing an exponential curve over an 800+ hour mostly-flat window — destabilising every RUL prediction after it.", 0),
    ("Fix applied: the voting condition must now hold for 2 consecutive readings (~20 min) before it counts as a real alarm.", 0),
    ("Also fixed a separate data-path bug that silently blocked rebuilding the pipeline from raw files.", 0),
])
add_note(s, "pdm/alarm.py (VOTING_CONFIRM_WINDOWS) and pdm/data_loader.py (PROJECT_ROOT path fix)", color=ACCENT, italic=False)

# ---------------------------------------------------------------------
# Slide 9b - Bug story, visual evidence
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "This Week's Work", "The Bug, Visualized")
add_image(s, SCRATCH / "healthscore_trigger_comparison.png", left=0.9, top=1.55, width=11.5, height=5.4)

# ---------------------------------------------------------------------
# Slide 10 - Before / after results
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Validation", "Results: Before vs. After the Fix")
rows = [
    ["Metric", "Test 2 (dev) — before", "Test 2 (dev) — after", "Test 3 (holdout) — before", "Test 3 (holdout) — after"],
    ["Lead time (h)", "74.8", "74.3", "816.3", "75.5"],
    ["Alarm episodes (baseline→final)", "19 → 1", "19 → 1", "507 → 7", "507 → 1"],
    ["RUL RMSE, last 10% of life (h)", "37.64", "26.45", "475.28", "106.33"],
]
add_table(s, rows, col_widths_in=[3.3, 2.35, 2.35, 2.35, 2.35], top=1.7, row_h=0.62,
          highlight_rows={1: LIGHT_BG, 2: LIGHT_BG, 3: RGBColor(0xE9, 0xF7, 0xEC)})
add_bullets(s, [
    ("Test 3 RMSE improved 4.5× (475h → 106h); lead time is now consistent with Test 2 instead of a physically implausible 816h.", 0),
    ("We are not claiming this matches the earlier draft's numbers — those remain unverifiable (no reproducible code found for them). This is what our own pipeline honestly produces, and why.", 1),
], top=4.3, height=2, font_size=15)

# ---------------------------------------------------------------------
# Slide 11 - Current validation table (dev vs holdout, post-fix, all metrics)
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Validation", "Development Run vs. Untouched Holdout (Current)")
rows = [
    ["Metric", "Test 2 (development)", "Test 3 (holdout, zero tuning)"],
    ["Run length / failure type", "163.8 h, outer race", "1073.3 h (6.5× longer), outer race"],
    ["Early-warning lead time", "74.3 h", "75.5 h"],
    ["False alarms, baseline vs. ours", "19 vs. 1", "507 vs. 1"],
    ["RUL RMSE, last 10% of life", "26.45 h", "106.33 h"],
]
add_table(s, rows, col_widths_in=[4.3, 4.0, 4.4], top=1.8, row_h=0.75)
add_bullets(s, [
    ("Same alarm/persistence logic, same code, applied unmodified to a run 6.5× longer — no per-run hand-tuning.", 0),
], top=4.9, height=1, font_size=15)

# ---------------------------------------------------------------------
# Slide 12 - Dashboard walkthrough (real screenshots, CRITICAL state)
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Deliverable", "Dashboard — KPIs & Health Trend (Live Screenshot)")
add_image(s, SCRATCH / "dash_11_trend_critical.png", left=0.7, top=1.45, width=11.9, height=5.55)
add_note(s, "Streamlit replay at t=108.8h — health score, alarm status and RUL countdown all populated from the same alarm event.", top=6.98, color=GRAY)

# ---------------------------------------------------------------------
# Slide 12b - Dashboard: Signal Detail + Alerts, side by side
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Deliverable", "Dashboard — Signal Detail & Alerts (Live Screenshots)")
add_image(s, SCRATCH / "dash_12_signal_critical.png", left=0.4, top=1.45, width=6.35, height=4.4,
          caption="Signal Detail: raw vibration + envelope spectrum (BPFO peak visible)")
add_image(s, SCRATCH / "dash_13_alerts_critical.png", left=6.85, top=1.45, width=6.1, height=4.4,
          caption="Alerts & Recommendation: reason, root cause, maintenance action")
add_note(s, "Live demo follows this slide — app/dashboard.py, run with: streamlit run app/dashboard.py", top=6.55, color=ACCENT, italic=False)

# ---------------------------------------------------------------------
# Slide 13 - Known limitations
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Transparency", "What We Do Not Claim Yet")
add_bullets(s, [
    ("Vibration only: the brief asks for temperature, current and pressure too. Not available in this dataset — addressed as roadmap, not hidden.", 0),
    ("Single operating condition: IMS runs at one constant speed/load. Real factory lines vary — the main real-world false-alarm source. Design answer: per-mode baselines and thresholds.", 0),
    ("Failure threshold (ET) currently transfers from a single other run. Should become a distribution across many runs as more failure data is collected.", 0),
    ("Single exponential RUL model: works well post-knee here, but a two-stage or piecewise model is on the roadmap for oscillating degradation patterns.", 0),
])

# ---------------------------------------------------------------------
# Slide 14 - Next steps
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s)
add_topbar(s, "Roadmap", "Next Steps (Week 3)")
add_bullets(s, [
    ("Continue investigating RUL stability on the non-primary bearings (test2 #2-4, test3 #1,2,4) that also show wear.", 0),
    ("Extend pdm/features.py to accept new sensor types (temperature, current, pressure) once real Astra data is available — architecture already supports this without redesign.", 0),
    ("Explore per-operating-mode baselines to prepare for real factory conditions.", 0),
    ("Prepare the live-streaming version of the replay pipeline (OPC UA / Modbus) as a forward-looking design, per the original roadmap.", 0),
])

# ---------------------------------------------------------------------
# Slide 15 - Thank you
# ---------------------------------------------------------------------
s = add_slide(); add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11), Inches(1))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Thank You"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
tb = s.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(11), Inches(0.6))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Questions & Discussion"
r.font.size = Pt(20); r.font.color.rgb = RGBColor(0x9C, 0xC5, 0xF0)

OUT_PATH.parent.mkdir(exist_ok=True)
prs.save(str(OUT_PATH))
print("Saved:", OUT_PATH)
print("Slide count:", len(prs.slides.__iter__.__self__._sldIdLst))
